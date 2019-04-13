import nltk
import collections
from nltk import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from sklearn.cluster import KMeans
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from pptx import Presentation
import os
import pickle

def prcoss(tokens):
    count = nltk.defaultdict(int)
    for word in tokens:
        count[word]+=1
    return count

def cos_sim(a,b):
    dot_product=np.dot(a,b)
    norm_a=np.linalg.norm(a)
    norm_b=np.linalg.norm(b)
    return dot_product/(norm_a*norm_b)

def getSimilarity(dict1, dict2):
    all_words_list=[]
    for key in dict1:
        all_words_list.append(key)
    for key in dict2:
        all_words_list.append(key)
    all_words_list_size=len(all_words_list)
    
    v1=np.zeros(all_words_list_size, dtype=np.int)
    v2=np.zeros(all_words_list_size, dtype=np.int)
    i=0
    for (key) in all_words_list:
        v1[i]=dict1.get(key, 0)
        v2[i]=dict2.get(key, 0)
        i=i+1
    return cos_sim(v1, v2)

def word_tokenizer(text):
    #tokenizes and stems the text
    tokens = word_tokenize(text)
    stemmer = PorterStemmer()
    tokens = [stemmer.stem(t) for t in tokens if t not in stopwords.words('english')]
    return tokens

def getPPTDetails(searchItem):
    path='/home/peter/Documents/myproject/media'
    directory=[x for x in os.listdir(path) if x.endswith(".pptx")]
    ppt={}
    for pptx_filename in directory:
        prs = Presentation(path+'/'+pptx_filename)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
        text_runs = []
        for slide in prs.slides:
            sh=''
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        sh=sh+' '+run.text
            if sh:
                text_runs.append(sh.split('.'))
        ppt[pptx_filename]=text_runs
    result=[]
    for keys in ppt:
        p=[]
        for i in range(0, len(ppt[keys])):
            for j in range(0, len(ppt[keys][i])):
                if getSimilarity(prcoss(word_tokenizer(ppt[keys][i][j])),prcoss(word_tokenizer(searchItem)))>0.1:
                    p.append([ppt[keys][i][j],i+1,j+1])
        if len(p):
            result.append([keys,p])
    return result

def ppttolist(file_name):
    path='media/'+file_name
    prs = Presentation(path)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    for slide in prs.slides:
        sh=''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    sh=sh+' '+run.text
        if sh:
            text_runs.append(sh.split('.'))
    fil='media/'+file_name.split('.')[0]+'.txt'
    with open(fil,'wb') as f:
        pickle.dump(text_runs,f)
    result=[]
    for i in range(0, len(text_runs)):
        p=[]
        for j in range(0, len(text_runs[i])):
            p.append(prcoss(word_tokenizer(text_runs[i][j])))
        if len(p):
            result.append(p)
    fil='media/'+file_name.split('.')[0]+'.pkl'
    with open(fil,'wb') as f:
        pickle.dump(result,f)

def cluster_sentences(sentences, nb_of_clusters=2):
    tfidf_vectorizer = TfidfVectorizer(tokenizer=word_tokenizer,
                       stop_words=stopwords.words('english'),
                       max_df=0.9,
                       min_df=0.1,
                       lowercase=True)
    #builds a tf-idf matrix for the sentences
    tfidf_matrix = tfidf_vectorizer.fit_transform(sentences)
    #print(tfidf_matrix)
    kmeans = KMeans(n_clusters=nb_of_clusters)
    kmeans.fit(tfidf_matrix)
    clusters = collections.defaultdict(list)
    for i, label in enumerate(kmeans.labels_):
        clusters[label].append(i)
    return dict(clusters)

def getPPTDetails1(searchItem):
    path='media/'
    directory=[x for x in os.listdir(path) if x.endswith(".pkl")]
    pkl={}
    for pkl_filename in directory:
        fil=path+pkl_filename
        with open(fil,'rb') as f:
            pkl[pkl_filename]=pickle.load(f)
    directory_txt=[x for x in os.listdir(path) if x.endswith(".txt")]
    ppt={}
    for txt_filename in directory_txt:
        fil=path+txt_filename
        with open(fil,'rb') as f:
            ppt[txt_filename]=pickle.load(f)
    result=[]
    for keys in pkl:
        key=keys.split('.')[0]+'.txt'
        p=[]
        for i in range(0, len(pkl[keys])):
            for j in range(0, len(pkl[keys][i])):
                k=getSimilarity(pkl[keys][i][j],prcoss(word_tokenizer(searchItem)))
                if k>0.1:
                    p.append([ppt[key][i][j],i+1,round(k*100,2)])
        if len(p):
            result.append([keys.split('.')[0],p])
    return result

def getPPTtext():
    path='media/'
    directory=[x for x in os.listdir(path) if x.endswith(".pptx")]
    ppt=[[], [], []]
    for pptx_filename in directory:
        prs = Presentation(path+pptx_filename)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
        text_runs = []
        j=1
        for slide in prs.slides:
            sh=''
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        sh=sh+' '+run.text
            if sh:
                text_runs+=sh.strip().split('.')
        j+=1
        if text_runs:
            ppt[0]+=text_runs
            for i in range(0, len(text_runs)):
                ppt[1].append(pptx_filename)
                ppt[2]+=str(j)
    return ppt
    
def clust(searchItem):
    data=getPPTtext()
    sentences=data[0]
    cl=[]
    nclusters=10
    clusters = cluster_sentences(sentences, nclusters)
    for cluster in range(nclusters):
        #print("cluster ",cluster+1,":")
        for i,sentence in enumerate(clusters[cluster]):
            if getSimilarity(prcoss(word_tokenizer(sentences[sentence])), prcoss(word_tokenizer(searchItem))) > 0:
                for i,sentence in enumerate(clusters[cluster]):
                    #print("\tsentence ",i+1,": ",sentences[sentence])
                    cl.append([sentences[sentence],data[1][sentence],str(data[2][sentence])])
                break

    return cl