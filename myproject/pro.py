import nltk
import collections
from nltk import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import numpy as np
from pptx import Presentation
import os


def prcoss(tokens):
    count = nltk.defaultdict(int)
    for word in tokens:
        count[word] += 1
    return count


def cos_sim(a, b):
    dot_product = np.dot(a, b)
    norm_a = np.linalg.norm(a)
    norm_b = np.linalg.norm(b)
    return dot_product / (norm_a * norm_b)


def getSimilarity(dict1, dict2):
    all_words_list = []
    for key in dict1:
        all_words_list.append(key)
    for key in dict2:
        all_words_list.append(key)
    all_words_list_size = len(all_words_list)

    v1 = np.zeros(all_words_list_size, dtype=np.int)
    v2 = np.zeros(all_words_list_size, dtype=np.int)
    i = 0
    for (key) in all_words_list:
        v1[i] = dict1.get(key, 0)
        v2[i] = dict2.get(key, 0)
        i = i + 1
    return cos_sim(v1, v2)


def word_tokenizer(text):
    # tokenizes and stems the text
    tokens = word_tokenize(text)
    stemmer = PorterStemmer()
    tokens = [stemmer.stem(t) for t in tokens if t not in stopwords.words('english')]
    return tokens


def getPPTDetails(searchItem):
    path = '/home/peter/Documents/myproject/media'
    directory = [x for x in os.listdir(path) if x.endswith(".pptx")]
    ppt = {}
    for pptx_filename in directory:
        prs = Presentation(path + '/' + pptx_filename)
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        text_runs = []
        for slide in prs.slides:
            sh = ''
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        sh = sh + ' ' + run.text
            if sh:
                text_runs.append(sh.split('.'))
        ppt[pptx_filename] = text_runs
    result = []
    for keys in ppt:
        p = []
        for i in range(0, len(ppt[keys])):
            for j in range(0, len(ppt[keys][i])):
                if getSimilarity(prcoss(word_tokenizer(ppt[keys][i][j])), prcoss(word_tokenizer(searchItem))) > 0.1:
                    p.append([ppt[keys][i][j], i + 1, j + 1])
        if len(p):
            result.append([keys, p])
    return result
