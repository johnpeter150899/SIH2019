#!/usr/bin/env python
# coding: utf-8

# In[1]:


from pptx import Presentation


# In[116]:


#from pptx import Presentation
prs = Presentation('F:\PPTs\deadlocks.pptx')
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
for slide in prs.slides:
    #sl=[]
    sh=''
    for shape in slide.shapes:
        #sh=[]
        #sh=''
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            #para=[]
            for run in paragraph.runs:
                #para.append(run.text)
                #sh.append(run.text)
                sh=sh+' '+run.text
        """if sh:
            sl.append(sh.split('.'))"""
    #text_runs.append(sl)
    if sh:
            text_runs.append(sh.split('.'))


# In[70]:


text_runs[0]


# In[79]:


import collections
from nltk import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from sklearn.cluster import KMeans
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np


# In[72]:


def word_tokenizer(text):
    #tokenizes and stems the text
    tokens = word_tokenize(text)
    stemmer = PorterStemmer()
    tokens = [stemmer.stem(t) for t in tokens if t not in stopwords.words('english')]
    return tokens


# In[73]:


def cluster_sentences(sentences, nb_of_clusters=5):
    tfidf_vectorizer = TfidfVectorizer(tokenizer=word_tokenizer,
                       stop_words=stopwords.words('english'),
                       max_df=0.9,
                       min_df=0.1,
                       lowercase=True)
    #builds a tf-idf matrix for the sentences
    tfidf_matrix = tfidf_vectorizer.fit_transform(sentences)
    #print(tfidf_matrix)
    kmeans = KMeans(n_clusters=nb_of_clusters)
    kmeans.fit(tfidf_matrix)
    clusters = collections.defaultdict(list)
    for i, label in enumerate(kmeans.labels_):
        clusters[label].append(i)
    return dict(clusters)


# In[74]:


import nltk
nltk.download('punkt')


# In[95]:


#for i in range(len(text_runs[1]))
sentences = text_runs[1]
nclusters= 3
#d={}
clusters = cluster_sentences(sentences, nclusters)
for cluster in range(nclusters):
    print("cluster ",cluster,":")
    for i,sentence in enumerate(clusters[cluster]):
        #d[]
        print("\tsentence ",i,": ",sentences[sentence])
clusters


# In[78]:


import os
os.listdir()


# In[91]:


def prcoss(tokens):
    count = nltk.defaultdict(int)
    for word in tokens:
        count[word]+=1
    return count

def cos_sim(a,b):
    dot_product=np.dot(a,b)
    norm_a=np.linalg.norm(a)
    norm_b=np.linalg.norm(b)
    return dot_product/(norm_a*norm_b)


# In[92]:


def getSimilarity(dict1, dict2):
    all_words_list=[]
    for key in dict1:
        all_words_list.append(key)
    for key in dict2:
        all_words_list.append(key)
    all_words_list_size=len(all_words_list)
    
    v1=np.zeros(all_words_list_size, dtype=np.int)
    v2=np.zeros(all_words_list_size, dtype=np.int)
    i=0
    for (key) in all_words_list:
        v1[i]=dict1.get(key, 0)
        v2[i]=dict2.get(key, 0)
        i=i+1
    return cos_sim(v1, v2)


# In[94]:


getSimilarity(prcoss(word_tokenizer(text_runs[1][0])),prcoss(word_tokenizer('implement using django')))


# In[84]:


text_runs[1]


# In[99]:


getSimilarity(prcoss(word_tokenizer("This presentation contains forward-looking statements and information that involve risks, uncertainties and assumptions.  Forward-looking statements are all statements that concern  plans, objectives, goals, strategies, future events or performance and underlying assumptions and other statements that are other than statements of historical fact, including, but not limited to, those that are identified by the use of words such as “anticipates”, “believes”, “estimates”, “expects”, “intends”, “plans”, “predicts”, “projects” and similar expressions. Risks and uncertainties that could affect us include, without limitation:")),prcoss(word_tokenizer('Gowtham Vasudevan John Peter')))


# In[125]:


textrun=np.array(text_runs)
sim=[]
for i in range(0, len(textrun)):
    ss=[]
    for j in range(0, len(text_runs[i])):
        ss.append(getSimilarity(prcoss(word_tokenizer(text_runs[i][j])),prcoss(word_tokenizer('super speciality therapies'))))
    sim.append(ss)
print(sim)


# In[ ]:





# In[ ]:




